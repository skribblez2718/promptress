"""
title: office_document_tool
author: sketch
version: 2.2.0
license: MIT
description: Create or modify Office documents (.docx, .pptx, .xlsx) and return them as downloadable attachments in Open WebUI chat.
requirements: python-docx,python-pptx,openpyxl,pydantic
"""

from __future__ import annotations

import asyncio
import base64
import io
import os
import uuid
from datetime import datetime
from typing import Any, Dict, List, Literal, Optional, Tuple, Union, cast

from pydantic import BaseModel, Field, ValidationError, field_validator

# Open WebUI internals - UPDATED IMPORTS for 0.5.x+
from open_webui.models.users import Users
from open_webui.models.files import Files, FileForm
from open_webui.storage.provider import Storage

# Office libs
from docx import Document as DocxDocument
from docx.shared import Inches, Pt
from docx.oxml.ns import qn
from pptx import Presentation
from pptx.util import Inches as PptInches
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from openpyxl import Workbook, load_workbook
from openpyxl.styles import PatternFill
from openpyxl.chart import BarChart, Reference
from openpyxl.formatting.rule import CellIsRule, ColorScaleRule


# ----------------------------
# Pydantic Schemas & Enums
# ----------------------------

FileType = Literal["docx", "pptx", "xlsx"]
OperationType = Literal["create", "modify"]


class ImageSpec(BaseModel):
    """Image to insert, by filename (from __files__) or bytes via base64."""

    name: Optional[str] = Field(
        default=None, description="File name of the image as uploaded/attached."
    )
    b64: Optional[str] = Field(
        default=None, description="Base64-encoded binary of the image."
    )
    width_inches: float = Field(
        default=2.0, ge=0.1, le=20.0, description="Image display width in inches."
    )


class ParagraphSpec(BaseModel):
    """Paragraph with optional inline formatting."""

    text: str = Field(..., description="Paragraph text.")
    bold: bool = Field(default=False)
    italic: bool = Field(default=False)
    underline: bool = Field(default=False)
    font_name: Optional[str] = Field(default=None)
    font_size_pt: Optional[float] = Field(default=None, ge=6.0, le=96.0)
    style: Optional[str] = Field(
        default=None,
        description="Word paragraph style name (e.g., 'Title', 'Heading 1').",
    )


class TableSpec(BaseModel):
    """Simple table from a 2D array of strings."""

    rows: List[List[str]] = Field(..., min_length=1, description="2D data for table.")
    style: Optional[str] = Field(
        default="Light List", description="Word table style if available."
    )


class FindReplaceSpec(BaseModel):
    """Naive find/replace across paragraph runs in Word."""

    find: str
    replace: str


class WordInstructions(BaseModel):
    """Operations for Word documents."""

    paragraphs: List[ParagraphSpec] = Field(default_factory=list)
    tables: List[TableSpec] = Field(default_factory=list)
    images: List[ImageSpec] = Field(default_factory=list)
    header_text: Optional[str] = None
    footer_text: Optional[str] = None
    find_replace: List[FindReplaceSpec] = Field(default_factory=list)
    default_style: Optional[str] = Field(default=None)


class SlideSpec(BaseModel):
    """Content for a PowerPoint slide."""

    title: Optional[str] = None
    bullets: List[str] = Field(default_factory=list)
    images: List[ImageSpec] = Field(default_factory=list)
    chart: Optional[Dict[str, Any]] = Field(
        default=None,
        description="Chart spec: {'type':'bar','categories':[...],'series':[{'name':..., 'values':[...]}]}",
    )


class PptInstructions(BaseModel):
    """Operations for PowerPoint presentations."""

    title: Optional[str] = None
    slides: List[SlideSpec] = Field(default_factory=list)
    note: Optional[str] = Field(
        default="Animations are limited by python-pptx and will be ignored if requested."
    )


class SheetSpec(BaseModel):
    """Excel sheet content and styling."""

    name: str = Field(..., description="Worksheet name.")
    data: List[List[Union[str, float, int, None]]] = Field(
        default_factory=list, description="2D grid data."
    )
    formulas: Dict[str, str] = Field(
        default_factory=dict, description="Cell formulas: {'B2': '=SUM(A1:A10)'}"
    )
    number_formats: Dict[str, str] = Field(
        default_factory=dict, description="Cell formats: {'B2':'0.00%'}"
    )
    column_widths: Dict[int, float] = Field(
        default_factory=dict, description="1-based col index -> width."
    )
    conditional_formatting: List[Dict[str, Any]] = Field(
        default_factory=list,
        description="List of {'type':'cellIs'|'colorScale', 'range':'A1:A10', ...}",
    )


class ExcelInstructions(BaseModel):
    """Operations for Excel workbooks."""

    sheets: List[SheetSpec] = Field(default_factory=list)
    chart: Optional[Dict[str, Any]] = Field(
        default=None,
        description="Workbook-level chart example on first sheet (bar) from a data range.",
    )


class OfficeToolParams(BaseModel):
    """
    Input schema for the office_document_tool.
    """

    file_type: FileType = Field(..., description="Office file type to create/modify.")
    operation: OperationType = Field(..., description="'create' or 'modify'.")
    instructions: Optional[
        Union[WordInstructions, PptInstructions, ExcelInstructions]
    ] = Field(
        default=None, description="Structured instructions; recommended over free-form."
    )
    raw_instructions: Optional[Dict[str, Any]] = Field(default=None)
    source_filename_hint: Optional[str] = Field(default=None)
    source_path: Optional[str] = Field(
        default=None, description="Existing file path for modify operations."
    )
    output_basename: Optional[str] = Field(
        default=None, description="Base name for the generated file (no extension)."
    )

    @field_validator("output_basename")
    @classmethod
    def _clean_basename(cls, v: Optional[str]) -> Optional[str]:
        if not v:
            return v
        safe = "".join(ch for ch in v if ch.isalnum() or ch in ("-", "_", " "))
        return safe.strip() or None


# ----------------------------
# Utilities
# ----------------------------


def _now_stamp() -> str:
    return datetime.now().strftime("%Y%m%d_%H%M%S")


def _choose_output_name(file_type: FileType, base: Optional[str]) -> str:
    root = base or f"document_{_now_stamp()}"
    ext = {"docx": ".docx", "pptx": ".pptx", "xlsx": ".xlsx"}[file_type]
    return f"{root}{ext}"


def _friendly_error(prefix: str, exc: Exception) -> str:
    return f"{prefix}: {type(exc).__name__} — {str(exc) or 'unexpected error'}"


def _get_content_type(ext: str) -> str:
    return {
        "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    }.get(ext, "application/octet-stream")


def _find_attached_file(
    expected_ext: str,
    files: Optional[List[Dict[str, Any]]] = None,
    hint: Optional[str] = None,
) -> Optional[Tuple[str, bytes]]:
    """
    Pick a file from __files__ by extension (and optional name hint).
    Accepts content either in 'content' (base64) or 'b64' keys or 'path' (not recommended).
    """
    if not files:
        return None

    # Try hint first
    if hint:
        for f in files:
            name = f.get("name") or f.get("filename") or ""
            if (
                name
                and name.lower() == hint.lower()
                and name.lower().endswith(expected_ext)
            ):
                data = None
                if "content" in f and isinstance(f["content"], str):
                    data = base64.b64decode(f["content"])
                elif "b64" in f and isinstance(f["b64"], str):
                    data = base64.b64decode(f["b64"])
                elif (
                    "path" in f
                    and isinstance(f["path"], str)
                    and os.path.isfile(f["path"])
                ):
                    with open(f["path"], "rb") as fh:
                        data = fh.read()
                if data:
                    return (name, data)

    # Otherwise first matching extension
    for f in files:
        name = f.get("name") or f.get("filename") or ""
        if name and name.lower().endswith(expected_ext):
            data = None
            if "content" in f and isinstance(f["content"], str):
                data = base64.b64decode(f["content"])
            elif "b64" in f and isinstance(f["b64"], str):
                data = base64.b64decode(f["b64"])
            elif (
                "path" in f and isinstance(f["path"], str) and os.path.isfile(f["path"])
            ):
                with open(f["path"], "rb") as fh:
                    data = fh.read()
            if data:
                return (name, data)

    return None


# ----------------------------
# Word (.docx) handlers
# ----------------------------


def _apply_word_instructions(doc: DocxDocument, instr: WordInstructions) -> None:
    """Apply Word operations (paragraphs, tables, images, header/footer, find/replace)."""
    if instr.default_style:
        try:
            doc.styles["Normal"].font.name = instr.default_style  # best-effort
        except Exception:
            pass

    for p in instr.paragraphs:
        para = doc.add_paragraph()
        run = para.add_run(p.text)
        run.bold = p.bold
        run.italic = p.italic
        run.underline = p.underline
        if p.font_name:
            run.font.name = p.font_name
            run._element.rPr.rFonts.set(qn("w:eastAsia"), p.font_name)
        if p.font_size_pt:
            run.font.size = Pt(p.font_size_pt)
        if p.style:
            para.style = p.style

    for t in instr.tables:
        rows = len(t.rows)
        cols = max(len(r) for r in t.rows)
        table = doc.add_table(rows=rows, cols=cols)
        if t.style:
            try:
                table.style = t.style
            except Exception:
                pass
        for i, row in enumerate(t.rows):
            for j, cell_text in enumerate(row):
                table.cell(i, j).text = str(cell_text)

    for im in instr.images:
        if im.b64:
            data = base64.b64decode(im.b64)
            doc.add_picture(io.BytesIO(data), width=Inches(im.width_inches))

    if instr.header_text:
        section = doc.sections[0]
        header = section.header
        header.paragraphs[0].text = instr.header_text

    if instr.footer_text:
        section = doc.sections[0]
        footer = section.footer
        footer.paragraphs[0].text = instr.footer_text

    # naive find/replace
    for fr in instr.find_replace:
        for para in doc.paragraphs:
            if fr.find in para.text:
                para.text = para.text.replace(fr.find, fr.replace)


def _create_docx(instr: WordInstructions) -> bytes:
    doc = DocxDocument()
    _apply_word_instructions(doc, instr)
    bio = io.BytesIO()
    doc.save(bio)
    return bio.getvalue()


def _modify_docx(existing: bytes, instr: WordInstructions) -> bytes:
    bio = io.BytesIO(existing)
    doc = DocxDocument(bio)
    _apply_word_instructions(doc, instr)
    out = io.BytesIO()
    doc.save(out)
    return out.getvalue()


# ----------------------------
# PowerPoint (.pptx) handlers
# ----------------------------


def _add_ppt_slide(prs: Presentation, slide_spec: SlideSpec) -> None:
    layout = prs.slide_layouts[1]  # Title + Content
    slide = prs.slides.add_slide(layout)

    if slide_spec.title:
        slide.shapes.title.text = slide_spec.title

    if slide_spec.bullets:
        body = slide.placeholders[1].text_frame
        body.clear()
        for i, b in enumerate(slide_spec.bullets):
            if i == 0:
                body.text = b
            else:
                p = body.add_paragraph()
                p.text = b
                p.level = 0

    for im in slide_spec.images:
        if im.b64:
            data = base64.b64decode(im.b64)
            slide.shapes.add_picture(
                io.BytesIO(data),
                PptInches(im.width_inches),
                PptInches(1.0),
                width=PptInches(im.width_inches),
            )

    # Simple bar chart example
    if slide_spec.chart:
        spec = slide_spec.chart
        ctype = (spec.get("type") or "bar").lower()
        if ctype == "bar":
            categories = spec.get("categories") or []
            series_list = spec.get("series") or []
            chart_data = ChartData()
            chart_data.categories = categories
            for s in series_list:
                chart_data.add_series(s.get("name", "Series"), s.get("values", []))
            x, y, cx, cy = PptInches(1), PptInches(3), PptInches(8), PptInches(3)
            slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
            )


def _create_pptx(instr: PptInstructions) -> bytes:
    prs = Presentation()
    if instr.title:
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = instr.title
        if slide.placeholders and len(slide.placeholders) > 1:
            slide.placeholders[1].text = ""
    for s in instr.slides:
        _add_ppt_slide(prs, s)
    bio = io.BytesIO()
    prs.save(bio)
    return bio.getvalue()


def _modify_pptx(existing: bytes, instr: PptInstructions) -> bytes:
    prs = Presentation(io.BytesIO(existing))
    for s in instr.slides:
        _add_ppt_slide(prs, s)
    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


# ----------------------------
# Excel (.xlsx) handlers
# ----------------------------


def _apply_sheet(ws, spec: SheetSpec) -> None:
    # Data grid
    for r_idx, row in enumerate(spec.data, start=1):
        for c_idx, value in enumerate(row, start=1):
            ws.cell(row=r_idx, column=c_idx, value=value)

    # Formulas
    for addr, formula in spec.formulas.items():
        ws[addr] = formula

    # Number formats
    for addr, fmt in spec.number_formats.items():
        ws[addr].number_format = fmt

    # Column widths (A=1, B=2, ...)
    for col_idx, width in spec.column_widths.items():
        ws.column_dimensions[chr(64 + col_idx)].width = float(width)

    # Conditional formatting
    for rule in spec.conditional_formatting:
        rtype = (rule.get("type") or "").lower()
        rng = rule.get("range") or ""
        if rtype == "cellis":
            op = rule.get("operator", "greaterThan")
            formula = rule.get("formula", "0")
            fill = PatternFill(
                start_color=rule.get("color", "FFC7CE"),
                end_color=rule.get("color", "FFC7CE"),
                fill_type="solid",
            )
            ws.conditional_formatting.add(
                rng, CellIsRule(operator=op, formula=[formula], fill=fill)
            )
        elif rtype == "colorscale":
            ws.conditional_formatting.add(
                rng,
                ColorScaleRule(
                    start_type="min",
                    start_color=rule.get("start_color", "63BE7B"),
                    mid_type="percentile",
                    mid_value=50,
                    mid_color=rule.get("mid_color", "FFEB84"),
                    end_type="max",
                    end_color=rule.get("end_color", "F8696B"),
                ),
            )


def _create_xlsx(instr: ExcelInstructions) -> bytes:
    wb = Workbook()
    if instr.sheets:
        wb.remove(wb.active)
    for s in instr.sheets:
        ws = wb.create_sheet(title=s.name[:31] or "Sheet")
        _apply_sheet(ws, s)

    # Simple chart on first sheet if requested
    if instr.chart and instr.sheets:
        ws = wb[instr.sheets[0].name[:31]]
        spec = instr.chart
        if (spec.get("type") or "bar").lower() == "bar":
            data_range = spec.get("data_range")  # e.g., "A1:D5"
            if data_range:
                min_col = ord(data_range.split(":")[0][0].upper()) - 64
                min_row = int(data_range.split(":")[0][1:])
                max_col = ord(data_range.split(":")[1][0].upper()) - 64
                max_row = int(data_range.split(":")[1][1:])
                data = Reference(
                    ws,
                    min_col=min_col,
                    min_row=min_row,
                    max_col=max_col,
                    max_row=max_row,
                )
                chart = BarChart()
                chart.add_data(data, titles_from_data=True)
                ws.add_chart(chart, "G2")

    bio = io.BytesIO()
    wb.save(bio)
    return bio.getvalue()


def _modify_xlsx(existing: bytes, instr: ExcelInstructions) -> bytes:
    wb = load_workbook(io.BytesIO(existing))
    for s in instr.sheets:
        ws = (
            wb[s.name]
            if s.name in wb.sheetnames
            else wb.create_sheet(title=s.name[:31])
        )
        _apply_sheet(ws, s)
    bio = io.BytesIO()
    wb.save(bio)
    return bio.getvalue()


# ----------------------------
# File Upload Helper (FIXED - correct FileForm structure)
# ----------------------------


def _upload_generated_file(
    file_bytes: bytes,
    filename: str,
    content_type: str,
    user_id: str,
    user_email: str,
    user_name: str,
) -> str:
    """
    Upload a generated file using Open WebUI's Storage provider and Files model.
    Returns the file ID for use in event emitter and URLs.

    This matches the exact structure used in open_webui/routers/files.py
    """
    # Generate unique file ID
    file_id = str(uuid.uuid4())

    # Create a unique storage filename to avoid collisions (matches router pattern)
    safe_filename = f"{file_id}_{filename}"

    # Build the tags dictionary (REQUIRED for S3StorageProvider)
    tags = {
        "OpenWebUI-User-Email": user_email,
        "OpenWebUI-User-Id": user_id,
        "OpenWebUI-User-Name": user_name,
        "OpenWebUI-File-Id": file_id,
    }

    # Upload to configured storage backend (local filesystem or S3)
    file_stream = io.BytesIO(file_bytes)
    contents, file_path = Storage.upload_file(file_stream, safe_filename, tags)

    # Create the file record in the database
    # IMPORTANT: FileForm structure must match what routers/files.py uses exactly
    file_form = FileForm(
        **{
            "id": file_id,
            "filename": filename,
            "path": file_path,  # TOP-LEVEL REQUIRED FIELD - not inside meta!
            "data": {},  # Required field, can be empty for generated files
            "meta": {
                "name": filename,
                "content_type": content_type,
                "size": len(file_bytes),
                "data": {"source": "office_document_tool"},
            },
        }
    )

    new_file = Files.insert_new_file(user_id, file_form)

    return new_file.id


# ----------------------------
# Main Tools class
# ----------------------------


class Tools:
    """
    Open WebUI tool entrypoint. Expose `office_document_tool(...)` to the LLM.

    This tool:
      - Validates inputs via Pydantic schemas
      - Creates / modifies Office docs
      - Uploads the file using Storage provider and Files model (0.5.x+ compatible)
      - Emits a 'files' event with the returned file id
      - Returns a friendly status string including the download URL
    """

    class Valves(BaseModel):
        allow_server_paths: bool = Field(
            default=False,
            description="If True, allow reading from source_path on the server (admin only).",
        )
        show_progress: bool = Field(
            default=True,
            description="If True, emit status events during generation/upload.",
        )
        min_progress_delay_ms: int = Field(
            default=700,
            description="Minimal delay to keep status visible (milliseconds).",
        )
        open_webui_url: str = Field(
            default="http://localhost:8080/",
            description="Base URL to build file download links.",
        )

    def __init__(self):
        self.valves = self.Valves()
        self.citation = False
        self.file_handler = True  # Prevents default RAG processing of generated files

    async def _emit_status(
        self, __event_emitter__, text: str, done: bool = False
    ) -> None:
        if __event_emitter__ and self.valves.show_progress:
            await __event_emitter__(
                {
                    "type": "status",
                    "data": {"description": text, "done": done, "hidden": False},
                }
            )

    async def office_document_tool(
        self,
        file_type: FileType,
        operation: OperationType,
        instructions: Optional[Dict[str, Any]] = None,
        raw_instructions: Optional[Dict[str, Any]] = None,
        source_filename_hint: Optional[str] = None,
        source_path: Optional[str] = None,
        output_basename: Optional[str] = None,
        __files__: Optional[List[Dict[str, Any]]] = None,
        __event_emitter__=None,
        __user__: Optional[Dict[str, Any]] = None,
    ) -> str:
        """
        Create or modify Office documents (.docx, .pptx, .xlsx) and attach them to the chat.

        Parameters
        ----------
        file_type : Literal["docx","pptx","xlsx"]
        operation : Literal["create","modify"]
        instructions : dict | None
            Structured instruction payload shaped like WordInstructions, PptInstructions, or ExcelInstructions.
        raw_instructions : dict | None
            Fallback raw instructions; used if `instructions` is omitted or partially structured.
        source_filename_hint : str | None
            Name hint for which attached file to modify (if multiple are present).
        source_path : str | None
            Server-side file path to an existing doc (admins only, requires valves.allow_server_paths = True).
        output_basename : str | None
            Desired base name for the generated file (extension will be added automatically).
        __files__ : list[dict] | None
            Files attached by the user to this message (Open WebUI provides these).
        __event_emitter__ : callable
            Event emitter provided by Open WebUI. We'll emit status and files events.
        __user__ : dict | None
            Current user (used to upload via OWUI file router).

        Returns
        -------
        str
            A short, user-friendly status string. The actual file(s) appear as attachments in the chat and as a link.
        """

        # Visible progress in the chat UI
        await self._emit_status(__event_emitter__, "Generating document…", done=False)
        if self.valves.min_progress_delay_ms > 0:
            await asyncio.sleep(self.valves.min_progress_delay_ms / 1000.0)

        try:
            # Build validated params (coerce instructions by file type)
            parsed = OfficeToolParams(
                file_type=file_type,
                operation=operation,
                instructions=None,
                raw_instructions=raw_instructions or instructions or {},
                source_filename_hint=source_filename_hint,
                source_path=source_path,
                output_basename=output_basename,
            )

            # Coerce instructions
            payload = parsed.raw_instructions or {}
            if parsed.file_type == "docx":
                instr_obj: Union[
                    WordInstructions, PptInstructions, ExcelInstructions
                ] = (WordInstructions(**payload) if payload else WordInstructions())
            elif parsed.file_type == "pptx":
                instr_obj = PptInstructions(**payload) if payload else PptInstructions()
            else:
                instr_obj = (
                    ExcelInstructions(**payload) if payload else ExcelInstructions()
                )

            # Resolve any image name -> bytes from __files__
            def resolve_images(images: List[ImageSpec]) -> List[ImageSpec]:
                if not images:
                    return images
                name_to_b64: Dict[str, str] = {}
                if __files__:
                    for f in __files__:
                        name = f.get("name") or f.get("filename")
                        if not name:
                            continue
                        if "content" in f and isinstance(f["content"], str):
                            try:
                                base64.b64decode(f["content"])
                                name_to_b64[name] = f["content"]
                            except Exception:
                                pass
                out: List[ImageSpec] = []
                for im in images:
                    if not im.b64 and im.name and im.name in name_to_b64:
                        out.append(
                            ImageSpec(
                                name=im.name,
                                b64=name_to_b64[im.name],
                                width_inches=im.width_inches,
                            )
                        )
                    else:
                        out.append(im)
                return out

            if isinstance(instr_obj, WordInstructions):
                instr_obj.images = resolve_images(instr_obj.images)
            elif isinstance(instr_obj, PptInstructions):
                for s in instr_obj.slides:
                    s.images = resolve_images(s.images)

            # Create or modify
            expected_ext = f".{parsed.file_type}"
            output_name = _choose_output_name(parsed.file_type, parsed.output_basename)

            if parsed.operation == "create":
                if parsed.file_type == "docx":
                    data_out = _create_docx(cast(WordInstructions, instr_obj))
                elif parsed.file_type == "pptx":
                    data_out = _create_pptx(cast(PptInstructions, instr_obj))
                else:
                    data_out = _create_xlsx(cast(ExcelInstructions, instr_obj))
            else:
                existing_bytes: Optional[bytes] = None
                pick = _find_attached_file(
                    expected_ext, __files__ or [], parsed.source_filename_hint
                )
                if pick:
                    _, existing_bytes = pick

                if not existing_bytes and parsed.source_path:
                    if self.valves.allow_server_paths and os.path.isfile(
                        parsed.source_path
                    ):
                        with open(parsed.source_path, "rb") as fh:
                            existing_bytes = fh.read()
                    else:
                        raise PermissionError(
                            "Server path access is disabled. Attach the file to the chat or enable allow_server_paths."
                        )
                if not existing_bytes:
                    raise FileNotFoundError(
                        f"No existing {expected_ext} file found to modify. Attach a file or provide a valid source_path."
                    )

                if parsed.file_type == "docx":
                    data_out = _modify_docx(
                        existing_bytes, cast(WordInstructions, instr_obj)
                    )
                elif parsed.file_type == "pptx":
                    data_out = _modify_pptx(
                        existing_bytes, cast(PptInstructions, instr_obj)
                    )
                else:
                    data_out = _modify_xlsx(
                        existing_bytes, cast(ExcelInstructions, instr_obj)
                    )

            # Upload using Storage provider + Files model (0.5.x+ compatible)
            await self._emit_status(
                __event_emitter__, "Uploading attachment…", done=False
            )

            # Resolve user ID from __user__ dict
            user_id = __user__.get("id") if isinstance(__user__, dict) else None
            if not user_id:
                raise RuntimeError("No user context was provided for upload.")

            # Get full user object to access email and name for tags
            user_obj = Users.get_user_by_id(user_id)
            if not user_obj:
                raise RuntimeError(f"User not found with ID: {user_id}")

            # Get content type
            content_type = _get_content_type(parsed.file_type)

            # Upload file using the fixed helper function (with tags and correct FileForm)
            file_id = _upload_generated_file(
                file_bytes=data_out,
                filename=output_name,
                content_type=content_type,
                user_id=user_id,
                user_email=user_obj.email or "",
                user_name=user_obj.name or "",
            )

            # Build file URL
            base_url = self.valves.open_webui_url.strip("/")
            file_url = f"{base_url}/api/v1/files/{file_id}/content"

            # Emit file attachment event (works in both Default and Native modes)
            if __event_emitter__:
                await __event_emitter__(
                    {
                        "type": "files",  # Use "files" for compatibility with both modes
                        "data": {
                            "files": [
                                {
                                    "type": "file",
                                    "id": file_id,
                                    "name": output_name,
                                    "url": f"/api/v1/files/{file_id}/content",
                                }
                            ]
                        },
                    }
                )
                await self._emit_status(__event_emitter__, "Done", done=True)

            # Final user-visible message below the attachment(s) including a direct link
            return (
                f"{parsed.operation.capitalize()}d {parsed.file_type.upper()} — "
                f"**{output_name}** is ready: [{output_name}]({file_url})"
            )

        except ValidationError as ve:
            if __event_emitter__:
                await __event_emitter__(
                    {
                        "type": "notification",
                        "data": {
                            "type": "error",
                            "content": f"Invalid inputs: {str(ve)}",
                        },
                    }
                )
            return _friendly_error("Input validation failed", ve)

        except PermissionError as pe:
            if __event_emitter__:
                await __event_emitter__(
                    {
                        "type": "notification",
                        "data": {"type": "warning", "content": str(pe)},
                    }
                )
            return _friendly_error("Permission error", pe)

        except FileNotFoundError as fe:
            if __event_emitter__:
                await __event_emitter__(
                    {
                        "type": "notification",
                        "data": {
                            "type": "warning",
                            "content": "No source file was found to modify.",
                        },
                    }
                )
            return _friendly_error("Missing source file", fe)

        except Exception as e:
            # Surface the actual error for debugging
            import traceback

            error_details = traceback.format_exc()
            if __event_emitter__:
                await __event_emitter__(
                    {
                        "type": "notification",
                        "data": {
                            "type": "error",
                            "content": f"Document processing failed: {str(e)}",
                        },
                    }
                )
            # Return more detailed error for debugging
            return f"Document processing failed: {type(e).__name__} — {str(e)}\n\nDetails:\n{error_details}"
